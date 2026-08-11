from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = "c:/Projects/AI-OPPORTUNIY-HUB/docs/AIHUB_BRIEF_EJECUTIVO_OPERATIVO.pptx"


def rgb(hex_code: str) -> RGBColor:
    hex_code = hex_code.lstrip("#")
    return RGBColor(int(hex_code[0:2], 16), int(hex_code[2:4], 16), int(hex_code[4:6], 16))


PALETTE = {
    "bg": rgb("F7F3EB"),
    "ink": rgb("1F2937"),
    "muted": rgb("5F6B7A"),
    "brand": rgb("0F766E"),
    "brand2": rgb("B45309"),
    "line": rgb("E6DDCF"),
    "white": rgb("FFFFFF"),
    "ok": rgb("166534"),
    "info": rgb("0C4A6E"),
}


def set_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = PALETTE["bg"]
    bg.line.fill.background()

    accent_left = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.4), Inches(-1.4), Inches(3.2), Inches(3.2))
    accent_left.fill.solid()
    accent_left.fill.fore_color.rgb = rgb("DCEFE9")
    accent_left.line.fill.background()

    accent_right = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.5), Inches(-1.3), Inches(3), Inches(3))
    accent_right.fill.solid()
    accent_right.fill.fore_color.rgb = rgb("F3E2CC")
    accent_right.line.fill.background()


def add_title(slide, title, subtitle=None):
    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.9), Inches(1.2))
    tf = tbox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Segoe UI"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = PALETTE["ink"]

    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.85), Inches(1.45), Inches(11.6), Inches(0.7))
        stf = sbox.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Segoe UI"
        sp.font.size = Pt(16)
        sp.font.color.rgb = PALETTE["muted"]


def add_card(slide, x, y, w, h, title, body, title_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PALETTE["white"]
    shape.line.color.rgb = PALETTE["line"]

    tf = shape.text_frame
    tf.clear()

    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(17)
    p1.font.bold = True
    p1.font.color.rgb = title_color or PALETTE["brand"]

    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(13)
    p2.font.color.rgb = PALETTE["muted"]
    p2.space_before = Pt(8)


def add_footer(slide, text="AI Opportunity Hub | Uso ejecutivo-operativo"):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.8), Inches(0.25))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Segoe UI"
    p.font.size = Pt(10)
    p.font.color.rgb = PALETTE["muted"]


def add_bullet_block(slide, x, y, w, h, heading, bullets):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = PALETTE["white"]
    box.line.color.rgb = PALETTE["line"]

    tf = box.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = heading
    p0.font.name = "Segoe UI"
    p0.font.size = Pt(18)
    p0.font.bold = True
    p0.font.color.rgb = PALETTE["ink"]

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.color.rgb = PALETTE["muted"]
        p.space_before = Pt(6)


def add_arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = PALETTE["brand2"]
    conn.line.width = Pt(2)


def build_deck():
    prs = Presentation()
    blank = prs.slide_layouts[6]

    # Slide 1: Portada
    s1 = prs.slides.add_slide(blank)
    set_bg(s1)
    add_title(
        s1,
        "AI Opportunity Hub",
        "Brief comercial ejecutivo-operativo para socializacion con clientes",
    )
    add_card(
        s1,
        Inches(0.9), Inches(2.1), Inches(11.6), Inches(2.35),
        "Propuesta de valor",
        "Sistema para transformar ideas de IA en iniciativas priorizadas y ejecutables, "
        "con validacion de negocio, filtro tecnico y paquete inicial de arquitectura.",
        PALETTE["brand"],
    )
    add_bullet_block(
        s1,
        Inches(0.9), Inches(4.7), Inches(11.6), Inches(1.8),
        "Mensajes clave para C-level",
        [
            "Decidir mas rapido donde invertir en IA",
            "Reducir riesgo de iniciativas sin factibilidad",
            "Mejorar conversion de ideas a ejecucion",
        ],
    )
    add_footer(s1)

    # Slide 2: Resumen ejecutivo
    s2 = prs.slides.add_slide(blank)
    set_bg(s2)
    add_title(s2, "Resumen ejecutivo", "Que es, que problema ataca y que resultado habilita")
    add_card(
        s2,
        Inches(0.9), Inches(2.0), Inches(5.65), Inches(2.0),
        "Que es",
        "Una capacidad de gobierno del pipeline de oportunidades IA: intake, evaluacion, "
        "priorizacion y salida accionable.",
    )
    add_card(
        s2,
        Inches(6.7), Inches(2.0), Inches(5.65), Inches(2.0),
        "Que resultado habilita",
        "Portafolio de iniciativas con trazabilidad de decisiones y menor friccion entre "
        "negocio, tecnologia y riesgo.",
    )
    add_bullet_block(
        s2,
        Inches(0.9), Inches(4.25), Inches(11.45), Inches(2.2),
        "Impacto esperado en 90 dias",
        [
            "-35% a -55% en tiempo de evaluacion inicial",
            "+20% a +40% en precision de priorizacion",
            "Criterios unificados para comites de decision",
        ],
    )
    add_footer(s2)

    # Slide 3: Problema
    s3 = prs.slides.add_slide(blank)
    set_bg(s3)
    add_title(s3, "Problema que resuelve", "Cuellos de botella tipicos en programas de IA")
    problems = [
        ("Backlog desordenado", "Muchas ideas, pocas con criterio uniforme de impacto y factibilidad."),
        ("Discovery costoso", "Se invierten semanas en analisis inicial sin salida comparable."),
        ("Decision fragmentada", "Negocio, tecnologia y riesgo evalian por separado y tarde."),
        ("Sin trazabilidad", "No hay claridad de por que una idea avanza o se rechaza."),
    ]
    x0 = [Inches(0.9), Inches(6.7)]
    y0 = [Inches(2.0), Inches(4.45)]
    idx = 0
    for y in y0:
        for x in x0:
            title, body = problems[idx]
            add_card(s3, x, y, Inches(5.65), Inches(2.05), title, body, PALETTE["brand2"])
            idx += 1
    add_footer(s3)

    # Slide 4: Objetivo
    s4 = prs.slides.add_slide(blank)
    set_bg(s4)
    add_title(s4, "Objetivo de la herramienta", "Estrategico y operativo")
    add_card(
        s4,
        Inches(0.9), Inches(2.1), Inches(5.65), Inches(2.6),
        "Objetivo estrategico",
        "Convertir el pipeline de ideas en una capacidad repetible de negocio para asignar "
        "inversion de IA con mayor certeza y menor riesgo.",
        PALETTE["info"],
    )
    add_card(
        s4,
        Inches(6.7), Inches(2.1), Inches(5.65), Inches(2.6),
        "Objetivo operativo",
        "Estandarizar intake, validacion por contexto de tenant, filtro tecnico y salida "
        "accionable para comites y equipos de entrega.",
        PALETTE["ok"],
    )
    add_bullet_block(
        s4,
        Inches(0.9), Inches(5.0), Inches(11.45), Inches(1.55),
        "Principio de diseno",
        ["Decisiones basadas en evidencia temprana, no en percepciones aisladas."],
    )
    add_footer(s4)

    # Slide 5: Diagrama de valor 1
    s5 = prs.slides.add_slide(blank)
    set_bg(s5)
    add_title(s5, "Diagrama de valor #1", "Cadena operativa de AI Opportunity Hub")

    labels = [
        ("Entrada", "Idea intake estructurado"),
        ("Filtro 1", "Validacion de negocio"),
        ("Filtro 2", "Validacion tecnica inicial"),
        ("Salida", "Architecture Package"),
    ]
    x_positions = [Inches(0.9), Inches(3.9), Inches(6.9), Inches(9.9)]
    for i, (tag, title) in enumerate(labels):
        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_positions[i], Inches(2.7), Inches(2.5), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = PALETTE["white"]
        card.line.color.rgb = PALETTE["line"]

        tf = card.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = tag
        p0.font.name = "Segoe UI"
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = PALETTE["brand"]

        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = PALETTE["ink"]
        p1.space_before = Pt(6)

    add_arrow(s5, Inches(3.4), Inches(3.8), Inches(3.85), Inches(3.8))
    add_arrow(s5, Inches(6.4), Inches(3.8), Inches(6.85), Inches(3.8))
    add_arrow(s5, Inches(9.4), Inches(3.8), Inches(9.85), Inches(3.8))

    add_bullet_block(
        s5,
        Inches(0.9), Inches(5.25), Inches(11.45), Inches(1.4),
        "Resultado operativo",
        ["Cada idea termina con recomendacion explicita: avanzar, ajustar o descartar."],
    )
    add_footer(s5)

    # Slide 6: KPI de valor
    s6 = prs.slides.add_slide(blank)
    set_bg(s6)
    add_title(s6, "Valor cuantificable", "Indicadores para patrocinadores y PMO")
    metrics = [
        ("-35% a -55%", "Reduccion esperada del tiempo de evaluacion inicial"),
        ("+20% a +40%", "Mejora esperada en precision de priorizacion"),
        ("1 marco comun", "Lenguaje unificado para negocio, TI y riesgo"),
        ("Trazabilidad", "Motivo de rechazo o avance por idea"),
    ]
    x0 = [Inches(0.9), Inches(6.7)]
    y0 = [Inches(2.0), Inches(4.2)]
    idx = 0
    for y in y0:
        for x in x0:
            title, body = metrics[idx]
            shape = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.65), Inches(1.95))
            shape.fill.solid()
            shape.fill.fore_color.rgb = PALETTE["white"]
            shape.line.color.rgb = PALETTE["line"]
            tf = shape.text_frame
            tf.clear()
            p1 = tf.paragraphs[0]
            p1.text = title
            p1.font.name = "Segoe UI"
            p1.font.size = Pt(24)
            p1.font.bold = True
            p1.font.color.rgb = PALETTE["brand2"]

            p2 = tf.add_paragraph()
            p2.text = body
            p2.font.name = "Segoe UI"
            p2.font.size = Pt(12)
            p2.font.color.rgb = PALETTE["muted"]
            p2.space_before = Pt(6)
            idx += 1
    add_footer(s6)

    # Slide 7: Diagrama de valor 2
    s7 = prs.slides.add_slide(blank)
    set_bg(s7)
    add_title(s7, "Diagrama de valor #2", "Escala de impacto de eficiencia a gobierno")

    levels = [
        ("Nivel 1", "Eficiencia", "Automatiza filtros iniciales y ahorra horas expertas."),
        ("Nivel 2", "Calidad de decision", "Compara ideas con criterios uniformes y menos sesgo."),
        ("Nivel 3", "Ejecucion", "Entrega paquetes iniciales para arrancar implementacion."),
        ("Nivel 4", "Gobierno", "Instala disciplina de portafolio y mejora por metrica."),
    ]

    for i, (lvl, name, desc) in enumerate(levels):
        y = Inches(2.0 + i * 1.25)
        x = Inches(1.0 + i * 0.5)
        w = Inches(9.8 - i * 0.45)
        shape = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(1.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PALETTE["white"]
        shape.line.color.rgb = PALETTE["line"]
        tf = shape.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = f"{lvl} | {name}"
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = PALETTE["brand"]

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(11)
        p2.font.color.rgb = PALETTE["muted"]

    add_footer(s7)

    # Slide 8: Ruta 90 dias
    s8 = prs.slides.add_slide(blank)
    set_bg(s8)
    add_title(s8, "Ruta de socializacion con clientes (90 dias)", "De workshop inicial a propuesta anual")

    phases = [
        ("Fase 0", "Semana 1", "Alineacion ejecutiva", "Workshop de alcance, KPIs y backlog inicial."),
        ("Fase 1", "Semanas 2-4", "Piloto controlado", "Evaluacion de 15-30 ideas reales y reporte quincenal."),
        ("Fase 2", "Mes 2", "Operacionalizacion", "Cadencia, RACI y comite de portafolio IA."),
        ("Fase 3", "Mes 3", "Escalamiento", "Roadmap anual por unidades y caso de negocio."),
    ]

    for i, (phase, when, title, desc) in enumerate(phases):
        y = Inches(2.0 + i * 1.25)
        box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), y, Inches(11.4), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = PALETTE["white"]
        box.line.color.rgb = PALETTE["line"]

        tf = box.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = f"{phase} ({when})"
        p0.font.name = "Segoe UI"
        p0.font.size = Pt(12)
        p0.font.bold = True
        p0.font.color.rgb = PALETTE["info"]

        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.name = "Segoe UI"
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = PALETTE["ink"]

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(11)
        p2.font.color.rgb = PALETTE["muted"]

    add_footer(s8)

    # Slide 9: Narrativa para cliente
    s9 = prs.slides.add_slide(blank)
    set_bg(s9)
    add_title(s9, "Narrativa comercial recomendada", "Mensajes diferenciados por audiencia")
    add_card(
        s9,
        Inches(0.9), Inches(2.0), Inches(5.65), Inches(3.0),
        "Mensaje ejecutivo (C-level)",
        "No es otro backlog de ideas. Es un sistema para decidir mejor y mas rapido donde invertir en IA, "
        "con evidencia y trazabilidad.",
        PALETTE["brand"],
    )
    add_card(
        s9,
        Inches(6.7), Inches(2.0), Inches(5.65), Inches(3.0),
        "Mensaje operativo (lideres de entrega)",
        "Reciben criterios estandarizados, motivo explicito de decision y paquetes de arquitectura "
        "listos para pasar de analisis a implementacion.",
        PALETTE["brand2"],
    )
    add_bullet_block(
        s9,
        Inches(0.9), Inches(5.2), Inches(11.45), Inches(1.4),
        "Enfoque de conversacion",
        ["Primero valor de negocio, luego riesgo controlado y por ultimo plan de ejecucion."],
    )
    add_footer(s9)

    # Slide 10: Oferta de piloto
    s10 = prs.slides.add_slide(blank)
    set_bg(s10)
    add_title(s10, "Oferta de piloto para cliente", "Alcance sugerido de entrada")
    add_bullet_block(
        s10,
        Inches(0.9), Inches(2.0), Inches(5.7), Inches(4.5),
        "Incluye",
        [
            "Kickoff ejecutivo y definicion de KPIs",
            "Carga y evaluacion de 15-30 ideas",
            "Top iniciativas priorizadas",
            "Recomendaciones de arquitectura inicial",
            "Roadmap de adopcion 6-12 meses",
        ],
    )
    add_bullet_block(
        s10,
        Inches(6.65), Inches(2.0), Inches(5.7), Inches(4.5),
        "Entregables",
        [
            "Informe ejecutivo de valor",
            "Tablero de decisiones con trazabilidad",
            "Playbook operativo de gobernanza",
            "Plan de escalamiento por unidad de negocio",
        ],
    )
    add_footer(s10)

    # Slide 11: Riesgos y mitigacion
    s11 = prs.slides.add_slide(blank)
    set_bg(s11)
    add_title(s11, "Riesgos de adopcion y mitigacion", "Marco de ejecucion responsable")
    risks = [
        ("Datos de baja calidad", "Definir criterios minimos de intake y ownership por dominio."),
        ("Expectativas de ROI poco realistas", "Alinear KPIs por fase y ventanas de valor (30/60/90 dias)."),
        ("Resistencia organizacional", "Patrocinio visible + comite transversal negocio/TI/riesgo."),
        ("Dependencia tecnologica", "Arquitectura modular con evolucion gradual hacia Foundry."),
    ]
    x0 = [Inches(0.9), Inches(6.7)]
    y0 = [Inches(2.0), Inches(4.3)]
    idx = 0
    for y in y0:
        for x in x0:
            title, body = risks[idx]
            add_card(s11, x, y, Inches(5.65), Inches(2.0), title, body, PALETTE["brand2"])
            idx += 1
    add_footer(s11)

    # Slide 12: Cierre
    s12 = prs.slides.add_slide(blank)
    set_bg(s12)
    add_title(s12, "Siguiente paso recomendado", "Cierre comercial")
    add_card(
        s12,
        Inches(0.9), Inches(2.05), Inches(11.45), Inches(2.4),
        "Propuesta de accion inmediata",
        "Agendar workshop ejecutivo de 90 minutos para definir objetivo de negocio, "
        "KPIs y universo inicial de ideas para el piloto.",
        PALETTE["ok"],
    )
    add_bullet_block(
        s12,
        Inches(0.9), Inches(4.7), Inches(11.45), Inches(1.85),
        "Compromisos de la semana 1",
        [
            "Cliente: sponsor + dueños de proceso + backlog inicial",
            "Equipo AI Opportunity Hub: metodologia, facilitacion y plan de evaluacion",
            "Salida: plan 90 dias con hitos de valor medibles",
        ],
    )
    add_footer(s12, "AI Opportunity Hub | Material comercial para sesiones de socializacion")

    prs.save(OUT_PATH)


if __name__ == "__main__":
    build_deck()
