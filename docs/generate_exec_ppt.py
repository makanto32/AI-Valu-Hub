from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = "docs/AIHUB_Brief_Ejecutivo_Operativo.pptx"


# Color palette
BG = RGBColor(243, 239, 231)
PRIMARY = RGBColor(15, 118, 110)
SECONDARY = RGBColor(180, 83, 9)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(95, 107, 122)
WHITE = RGBColor(255, 255, 255)
BORDER = RGBColor(230, 221, 207)
INFO = RGBColor(12, 74, 110)
OK = RGBColor(22, 101, 52)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color=BG):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.2), Inches(1.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = TEXT
    run.font.name = "Segoe UI"

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = MUTED
        p2.font.name = "Segoe UI"


def add_footer(slide, text="AI Opportunity Hub | Confidential"): 
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.font.name = "Segoe UI"


def add_card(slide, x, y, w, h, title, body, title_color=PRIMARY):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER

    t = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(w - 0.4), Inches(h - 0.25))
    tf = t.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = "Segoe UI"

    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(12)
    p2.font.color.rgb = MUTED
    p2.font.name = "Segoe UI"


def add_metric(slide, x, y, w, h, value, label):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BORDER

    t = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(w - 0.4), Inches(h - 0.25))
    tf = t.text_frame

    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT
    p.font.name = "Segoe UI"

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(12)
    p2.font.color.rgb = MUTED
    p2.font.name = "Segoe UI"


def add_flow_node(slide, x, y, w, h, tag, title, body):
    node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    node.fill.solid()
    node.fill.fore_color.rgb = WHITE
    node.line.color.rgb = BORDER

    tb = slide.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.1), Inches(w - 0.24), Inches(h - 0.2))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = tag
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.font.name = "Segoe UI"

    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = TEXT
    p2.font.name = "Segoe UI"

    p3 = tf.add_paragraph()
    p3.text = body
    p3.font.size = Pt(10)
    p3.font.color.rgb = MUTED
    p3.font.name = "Segoe UI"


def add_arrow(slide, x, y):
    a = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(0.5), Inches(0.6))
    a.fill.solid()
    a.fill.fore_color.rgb = SECONDARY
    a.line.fill.background()


# Slide 1: Cover
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(
    slide,
    "AI Opportunity Hub",
    "Brief ejecutivo y operativo para posicionamiento comercial",
)
add_card(
    slide,
    0.7,
    1.8,
    5.9,
    2.5,
    "Propuesta central",
    "Transformar ideas de IA en decisiones ejecutables, con trazabilidad y menor riesgo de inversion.",
)
add_card(
    slide,
    6.8,
    1.8,
    5.8,
    2.5,
    "Estado actual",
    "MVP1 implementado, MVP2 en progreso: validacion tecnica, architecture package y composicion de respuesta.",
    title_color=INFO,
)
add_metric(slide, 0.7, 4.6, 3.9, 1.6, "90 dias", "ruta propuesta para socializacion y adopcion")
add_metric(slide, 4.8, 4.6, 3.9, 1.6, "2 filtros", "negocio + tecnica para priorizacion")
add_metric(slide, 8.9, 4.6, 3.7, 1.6, "1 lenguaje", "ejecutivo, operativo y arquitectura")
add_footer(slide)

# Slide 2: Agenda
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Agenda de conversacion con cliente")
add_card(
    slide,
    0.8,
    1.5,
    12.0,
    5.0,
    "Estructura sugerida (45-60 min)",
    "1) Objetivo estrategico y operativo\n"
    "2) Problemas actuales de priorizacion IA\n"
    "3) Diagrama de valor de AI Opportunity Hub\n"
    "4) Potencial de escalamiento y gobierno\n"
    "5) Ruta de socializacion en 90 dias\n"
    "6) Acuerdos de siguiente paso",
)
add_footer(slide)

# Slide 3: Objectives
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Objetivo de la herramienta")
add_card(
    slide,
    0.8,
    1.5,
    6.1,
    4.8,
    "Objetivo estrategico",
    "Convertir el pipeline de ideas de IA en una capacidad repetible de negocio: detectar valor temprano, descartar rapido lo no viable y enfocar inversion en casos con impacto.",
)
add_card(
    slide,
    6.95,
    1.5,
    5.9,
    4.8,
    "Objetivo operativo",
    "Estandarizar intake, evaluacion por contexto de tenant, validacion tecnica inicial y salida accionable para patrocinadores, equipos de IA y arquitectura.",
    title_color=INFO,
)
add_footer(slide)

# Slide 4: Problem
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Problema que resuelve")
add_card(slide, 0.8, 1.5, 6.0, 2.2, "Backlog sin priorizacion", "Muchas ideas y poca claridad sobre retorno y factibilidad.", title_color=SECONDARY)
add_card(slide, 6.9, 1.5, 5.9, 2.2, "Discovery costoso", "Demasiado tiempo en analisis inicial sin artefactos comparables.", title_color=SECONDARY)
add_card(slide, 0.8, 3.9, 6.0, 2.2, "Decision fragmentada", "Negocio, tecnologia y riesgo evaluan en tiempos y formatos distintos.", title_color=SECONDARY)
add_card(slide, 6.9, 3.9, 5.9, 2.2, "Baja trazabilidad", "Poca capacidad para explicar por que una idea avanza o se descarta.", title_color=SECONDARY)
add_footer(slide)

# Slide 5: Value diagram 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Diagrama de valor #1 | Cadena operacional")
add_flow_node(slide, 0.7, 2.1, 2.8, 2.5, "ENTRADA", "Idea intake", "Captura uniforme de problema, objetivo y criterio de exito.")
add_arrow(slide, 3.6, 3.0)
add_flow_node(slide, 4.2, 2.1, 2.8, 2.5, "FILTRO 1", "Validacion negocio", "Prioriza por impacto esperado y alineacion estrategica.")
add_arrow(slide, 7.1, 3.0)
add_flow_node(slide, 7.7, 2.1, 2.8, 2.5, "FILTRO 2", "Validacion tecnica", "Evalua datos, integraciones, complejidad y riesgo.")
add_arrow(slide, 10.6, 3.0)
add_flow_node(slide, 11.2, 2.1, 1.5, 2.5, "SALIDA", "Paquete", "Architecture package y siguientes pasos.")
add_metric(slide, 0.8, 5.1, 4.0, 1.3, "-35% a -55%", "reduccion estimada en tiempo de evaluacion inicial")
add_metric(slide, 4.95, 5.1, 3.9, 1.3, "+20% a +40%", "mejor precision de priorizacion")
add_metric(slide, 9.0, 5.1, 3.8, 1.3, "Trazabilidad", "motivo de avance / rechazo por idea")
add_footer(slide)

# Slide 6: KPI and outcomes
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Resultados esperados y KPIs")
add_metric(slide, 0.8, 1.7, 2.9, 1.8, "T1", "Tiempo medio de evaluacion por idea")
add_metric(slide, 3.95, 1.7, 2.9, 1.8, "T2", "Tasa idea -> iniciativa priorizada")
add_metric(slide, 7.1, 1.7, 2.9, 1.8, "T3", "Tiempo de decision de comite")
add_metric(slide, 10.25, 1.7, 2.3, 1.8, "T4", "Calidad del backlog")
add_card(
    slide,
    0.8,
    3.8,
    12.0,
    2.6,
    "Lectura ejecutiva",
    "La plataforma no solo mejora eficiencia. Tambien aumenta la calidad de inversion en IA al elevar la consistencia de decision entre negocio, arquitectura y riesgo.",
    title_color=OK,
)
add_footer(slide)

# Slide 7: Value diagram 2 potential ladder
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Diagrama de valor #2 | Escala de potencial")
levels = [
    ("Nivel 1 | Eficiencia", "Automatiza filtros iniciales y libera capacidad analitica."),
    ("Nivel 2 | Calidad", "Comparacion objetiva de ideas bajo criterios comunes."),
    ("Nivel 3 | Ejecucion", "Salida con arquitectura inicial y siguientes pasos accionables."),
    ("Nivel 4 | Gobierno", "Disciplina de portafolio IA y mejora continua por metrica."),
]
y = 5.5
w = 11.7
for i, (ttl, body) in enumerate(levels):
    h = 0.9
    x = 0.8 + i * 0.5
    boxw = w - i * 1.0
    add_card(slide, x, y - i * 1.0, boxw, h, ttl, body, title_color=PRIMARY if i < 2 else INFO)
add_footer(slide)

# Slide 8: Operating model
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Modelo operativo propuesto")
add_card(slide, 0.8, 1.6, 4.0, 2.1, "Rol sponsor", "Define objetivos de valor y criterios de aprobacion.")
add_card(slide, 4.95, 1.6, 4.0, 2.1, "Rol negocio/ops", "Ingresa ideas, aporta contexto y valida prioridad.")
add_card(slide, 9.1, 1.6, 3.7, 2.1, "Rol arquitectura", "Valida factibilidad y propone blueprint inicial.")
add_card(slide, 0.8, 3.95, 12.0, 2.3, "Cadencia", "Ritmo quincenal para screening de ideas y comite mensual de decisiones. RACI claro, criterios estandarizados y trazabilidad de decisiones.")
add_footer(slide)

# Slide 9: 90-day roadmap
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Ruta de socializacion con clientes | 90 dias")
add_card(slide, 0.8, 1.7, 2.9, 4.8, "Fase 0\nSemana 1", "Alineacion ejecutiva\nWorkshop de objetivo, KPIs y backlog candidato.", title_color=INFO)
add_card(slide, 3.95, 1.7, 2.9, 4.8, "Fase 1\nSem 2-4", "Piloto controlado\n15-30 ideas reales con trazabilidad.", title_color=INFO)
add_card(slide, 7.1, 1.7, 2.9, 4.8, "Fase 2\nMes 2", "Operacionalizacion\nCadencia, RACI y comite de portafolio.", title_color=INFO)
add_card(slide, 10.25, 1.7, 2.6, 4.8, "Fase 3\nMes 3", "Escalamiento\nRoadmap anual y business case.", title_color=INFO)
add_footer(slide)

# Slide 10: Comercialization playbook
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Playbook comercial recomendado")
add_card(
    slide,
    0.8,
    1.6,
    6.0,
    4.8,
    "Narrativa para C-level",
    '"No es otro backlog de ideas. Es un sistema para decidir mejor y mas rapido donde invertir en IA, con evidencia y trazabilidad."',
    title_color=SECONDARY,
)
add_card(
    slide,
    6.95,
    1.6,
    5.85,
    4.8,
    "Narrativa para lideres operativos",
    '"Reciben un flujo claro de trabajo, motivos de rechazo explicitos y paquetes de arquitectura listos para mover iniciativas a ejecucion."',
    title_color=SECONDARY,
)
add_footer(slide)

# Slide 11: Next steps
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Siguiente paso sugerido")
add_card(
    slide,
    0.8,
    1.8,
    12.0,
    4.2,
    "Propuesta de accion inmediata",
    "1) Seleccionar cliente piloto y sponsor.\n"
    "2) Ejecutar workshop de alineacion (semana 1).\n"
    "3) Cargar backlog inicial y definir KPIs base.\n"
    "4) Revisar resultados al cierre de la semana 4 con decision de escalamiento.",
    title_color=OK,
)
add_metric(slide, 0.8, 6.2, 12.0, 0.9, "Entrega", "Deck listo para reuniones ejecutivas y operativas")
add_footer(slide)

prs.save(OUTPUT)
print(f"PPT generated: {OUTPUT}")
