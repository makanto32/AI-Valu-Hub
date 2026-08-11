from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUTPUT = "docs/AIHUB_Abordaje_Interno_Microsoft_FY.pptx"

BG = RGBColor(245, 247, 251)
TEXT = RGBColor(15, 23, 42)
MUTED = RGBColor(71, 85, 105)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(219, 227, 239)
BRAND = RGBColor(0, 120, 212)
BRAND2 = RGBColor(15, 118, 110)
ACCENT = RGBColor(180, 83, 9)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def add_title(slide, title, subtitle=""):
    box = slide.shapes.add_textbox(Inches(0.65), Inches(0.3), Inches(12.0), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.font.color.rgb = TEXT
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.name = "Segoe UI"
        p2.font.color.rgb = MUTED


def add_card(slide, x, y, w, h, title, body, color=BRAND):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = WHITE
    s.line.color.rgb = LINE

    b = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.15), Inches(w - 0.4), Inches(h - 0.2))
    tf = b.text_frame
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(16)
    p0.font.bold = True
    p0.font.color.rgb = color
    p0.font.name = "Segoe UI"

    p1 = tf.add_paragraph()
    p1.text = body
    p1.font.size = Pt(12)
    p1.font.color.rgb = MUTED
    p1.font.name = "Segoe UI"


def add_footer(slide):
    b = slide.shapes.add_textbox(Inches(0.65), Inches(7.03), Inches(12.0), Inches(0.25))
    p = b.text_frame.paragraphs[0]
    p.text = "AI Opportunity Hub | Socializacion interna Microsoft FY"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.font.name = "Segoe UI"


def add_arrow(slide, x, y):
    a = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(0.45), Inches(0.55))
    a.fill.solid()
    a.fill.fore_color.rgb = ACCENT
    a.line.fill.background()


# Slide 1
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_title(s, "AI Opportunity Hub", "Abordaje para socializacion interna con equipos Microsoft (proximo FY)")
add_card(s, 0.75, 1.8, 12.0, 2.3, "Objetivo", "Posicionar la propuesta como acelerador de pipeline calificado, decision rapida y ejecucion gobernada de iniciativas de IA.")
add_card(s, 0.75, 4.35, 3.9, 1.8, "Resultado esperado", "Inclusion en planes de cuenta FY")
add_card(s, 4.9, 4.35, 3.9, 1.8, "Horizonte", "Ruta 30-60-90 dias")
add_card(s, 9.05, 4.35, 3.7, 1.8, "Metrica norte", "Conversion idea -> iniciativa")
add_footer(s)

# Slide 2
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_title(s, "Mensaje central interno", "Como explicarlo en 30 segundos")
add_card(s, 0.75, 1.7, 12.0, 4.9, "Narrativa", "AI Opportunity Hub transforma conversaciones de IA en un sistema de decision medible: menos tiempo en discovery, mejor priorizacion y mas oportunidades convertidas en ejecucion sobre el stack Microsoft.", BRAND2)
add_footer(s)

# Slide 3
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_title(s, "Alineacion con estrategia FY", "Matriz de encaje")
add_card(s, 0.75, 1.7, 5.9, 2.2, "AI Transformation", "Genera backlog priorizado con criterios de valor y factibilidad.")
add_card(s, 6.85, 1.7, 5.9, 2.2, "Data + AI en Azure", "Conecta oportunidades a arquitectura y servicios objetivo.")
add_card(s, 0.75, 4.1, 5.9, 2.2, "Security + Governance", "Incorpora trazabilidad y rationale de decisiones.")
add_card(s, 6.85, 4.1, 5.9, 2.2, "Ejecucion Comercial", "Estandariza discovery y acelera paso a propuesta.")
add_footer(s)

# Slide 4
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_title(s, "Diagrama de valor interno", "Pipeline -> Velocity -> Attach -> Outcome")
add_card(s, 0.75, 2.4, 2.8, 2.1, "Pipeline", "Mas oportunidades calificadas por cuenta")
add_arrow(s, 3.65, 3.15)
add_card(s, 4.15, 2.4, 2.8, 2.1, "Velocity", "Decision mas rapida y menos retrabajo")
add_arrow(s, 7.05, 3.15)
add_card(s, 7.55, 2.4, 2.8, 2.1, "Attach", "Mayor vinculo a servicios de plataforma")
add_arrow(s, 10.45, 3.15)
add_card(s, 10.95, 2.4, 1.8, 2.1, "Outcome", "Win-rate")
add_card(s, 0.75, 5.0, 12.0, 1.5, "Mensaje clave", "No es una herramienta aislada: es un motor de decision para calidad de pipeline comercial y ejecucion tecnica.", BRAND2)
add_footer(s)

# Slide 5
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_title(s, "Stakeholders internos y mensaje por rol")
add_card(s, 0.75, 1.7, 3.9, 2.2, "Account Team", "Mejor pipeline y mas conversion a oportunidades reales.")
add_card(s, 4.85, 1.7, 3.9, 2.2, "Specialists Data + AI", "Mayor attach de servicios y casos mejor definidos.")
add_card(s, 8.95, 1.7, 3.8, 2.2, "Cloud Solution Architect", "Factibilidad temprana y menor retrabajo tecnico.")
add_card(s, 0.75, 4.1, 5.9, 2.2, "Customer Success / Delivery", "Entradas mas maduras para implementacion y adopcion.")
add_card(s, 6.85, 4.1, 5.9, 2.2, "Industry / Strategy Lead", "Oferta repetible por industria y escala en FY.")
add_footer(s)

# Slide 6
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_title(s, "KPIs para demostrar valor FY")
add_card(s, 0.75, 1.7, 3.9, 2.0, "KPI 1", "Pipeline calificado: ideas evaluadas vs priorizadas", ACCENT)
add_card(s, 4.85, 1.7, 3.9, 2.0, "KPI 2", "Time-to-decision por idea", ACCENT)
add_card(s, 8.95, 1.7, 3.8, 2.0, "KPI 3", "Attach de plataforma por iniciativa", ACCENT)
add_card(s, 0.75, 4.0, 3.9, 2.0, "KPI 4", "Conversion a propuesta/piloto", ACCENT)
add_card(s, 4.85, 4.0, 3.9, 2.0, "KPI 5", "Velocidad aprobacion -> plan", ACCENT)
add_card(s, 8.95, 4.0, 3.8, 2.0, "KPI 6", "Calidad de decision con trazabilidad", ACCENT)
add_footer(s)

# Slide 7
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_title(s, "Ruta de socializacion 30-60-90", "Plan operativo interno")
add_card(s, 0.75, 1.8, 3.9, 4.7, "Dia 0-30", "Alineacion y sponsor.\nValidar narrativa.\nSeleccionar 3-5 cuentas candidatas.", BRAND)
add_card(s, 4.85, 1.8, 3.9, 4.7, "Dia 31-60", "Prueba de valor.\nOpportunity shaping con backlog real.\nReporte quincenal de impacto.", BRAND2)
add_card(s, 8.95, 1.8, 3.8, 4.7, "Dia 61-90", "Escalamiento FY.\nIncluir resultados en QBR.\nPropuesta formal por industria/cuenta.", ACCENT)
add_footer(s)

# Slide 8
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_title(s, "Estructura de reunion interna (45 min)", "Formato sugerido")
add_card(s, 0.75, 1.7, 12.0, 4.9, "Agenda", "5 min: contexto FY y dolor actual\n10 min: diagrama de valor\n10 min: casos por cuenta\n10 min: KPIs y baseline\n10 min: acuerdos de ownership y siguientes pasos", BRAND2)
add_footer(s)

# Slide 9
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_title(s, "Objeciones internas y respuesta", "Preparacion para adopcion")
add_card(s, 0.75, 1.7, 5.9, 2.2, "Objecion 1", "Ya tenemos discovery.\nRespuesta: no reemplaza discovery, lo estandariza y lo acelera.")
add_card(s, 6.85, 1.7, 5.9, 2.2, "Objecion 2", "No hay tiempo para otro proceso.\nRespuesta: reduce ciclos y evita invertir en ideas no viables.")
add_card(s, 0.75, 4.1, 5.9, 2.2, "Objecion 3", "Valor dificil de medir.\nRespuesta: baseline y KPI quincenal desde semana 1.")
add_card(s, 6.85, 4.1, 5.9, 2.2, "Objecion 4", "Sin sponsor no escala.\nRespuesta: sponsor es hito de fase 0.")
add_footer(s)

# Slide 10
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_title(s, "Decision solicitada al equipo interno", "Call to action")
add_card(s, 0.75, 1.8, 12.0, 4.9, "Pedido concreto", "1) Aprobar piloto en 1-2 cuentas\n2) Nombrar sponsor interno\n3) Asignar owner tecnico-comercial\n4) Agendar workshop con cliente en 2 semanas\n5) Acordar 3 KPI de exito para 90 dias", BRAND)
add_footer(s)

prs.save(OUTPUT)
print(f"PPT generated: {OUTPUT}")
